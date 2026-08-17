"""Inventaire reproductible et prudent des ressources candidates de E:\\Documents.

Le script ne copie aucun contenu source. Il produit uniquement un manifeste de
métadonnées, des empreintes lorsque leur calcul reste borné et de courts
extraits destinés au classement. Toute qualification juridique reste humaine.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import re
import sqlite3
import zipfile
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from docx import Document  # type: ignore[import-not-found]
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-not-found]
from pypdf import PdfReader  # type: ignore[import-not-found]


DOCUMENT_EXTENSIONS = {".pdf", ".doc", ".docx", ".odt", ".rtf", ".ppt", ".pptx"}
TABLE_EXTENSIONS = {".csv", ".tsv", ".xls", ".xlsx", ".ods"}
RASTER_EXTENSIONS = {".tif", ".tiff", ".vrt", ".asc"}
VECTOR_EXTENSIONS = {".gpkg", ".geojson", ".kml", ".shp"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".webp"}
SHAPEFILE_SIDECARS = {".dbf", ".shx", ".prj", ".cpg", ".qix", ".sbn", ".sbx", ".fix"}

EXCLUDED_ROOTS = {
    "clès usb léandre": "archive personnelle hors périmètre",
    "codex": "sorties d'outillage, non-source scientifique",
    "autopsy_addon_modules-master": "code tiers sans rapport direct avec GSIE",
    "diskgenius professional 5.5.0.1488 win x64 multi + crack": "logiciel tiers et contenu à risque",
    "feedbackhub": "composant logiciel tiers",
    "honorsuite": "composant logiciel tiers",
    "ilovepdf": "composant logiciel tiers",
    "modèles office personnalisés": "modèles bureautiques génériques",
    "papier important": "documents personnels sensibles hors périmètre",
    "powershell": "outillage système",
    "windowspowershell": "outillage système",
    "volley ball meymac": "activité personnelle hors périmètre GSIE",
}

EXCLUDED_PATH_TOKENS = (
    "\\.git\\", "\\node_modules\\", "\\__pycache__\\", "\\build\\", "\\dist\\",
    "\\site-packages\\", "\\diagnostics\\stress_test_", "pdfelementportable", "\\app\\pdfelement\\",
    "convention de stage", "curriculum vitae", "\\cv.", "\\cv ", "contacts_forensic",
)

DOMAIN_PATTERNS: dict[str, tuple[str, ...]] = {
    "diagnostic_stationnel": (
        "diagnostic station", "stationnel", "station forest", "fiche terrain",
        "releve terrain", "relevé terrain", "placette", "ecologie forest",
    ),
    "pedologie": (
        "pedolog", "pédolog", "sol", "ph ", "rum", "reserve utile", "réserve utile",
        "texture", "hydromorph", "geolog", "géolog", "soil",
    ),
    "sylviculture": (
        "sylvic", "martelage", "futaie", "taillis", "coupe", "itineraire", "itinéraire",
        "hetraie", "hêtraie", "chenaie", "chênaie", "sapiniere", "sapinière",
    ),
    "dendrometrie": (
        "dendrom", "cubage", "surface terriere", "surface terrière", "volume", "diametre",
        "diamètre", "hauteur", "densite", "densité", "inventaire forest",
    ),
    "sante_pathologie": (
        "pathogene", "pathogène", "maladie", "champignon", "ravageur", "deperissement",
        "dépérissement", "sanitaire",
    ),
    "climat_hydrologie": (
        "climat", "pluvi", "temperature", "température", "secheresse", "sécheresse",
        "hydrolog", "zone humide", "inond", "meteo", "météo",
    ),
    "biodiversite": (
        "biodivers", "habitat", "flore", "faune", "botani", "espece", "espèce",
        "taxon", "naturaliste",
    ),
    "foret_gibier": (
        "gibier", "ongule", "ongulé", "cervide", "cervidé", "sylvo-cyneget",
        "sylvo-cynégét", "abroutissement",
    ),
    "economie_bois": (
        "prix du bois", "commercialisation", "estimation", "valeur de coupe", "exploitation",
        "cout", "coût", "rendement",
    ),
    "geospatial_teledetection": (
        "qgis", "lidar", "mnt", "mns", "mnh", "bd foret", "bd forêt", "orthophoto",
        "parcellaire", "geojson", "geopackage", "shapefile", "cartograph",
    ),
    "incendie_dfci": ("dfci", "incendie", "feu de foret", "feu de forêt", "combustible"),
    "benchmark_formation": (
        "rapport de stage", "placette eil", "chantier ecole", "chantier école",
        "diagnostic", "analyse de parcelle",
    ),
}

PRODUCER_PATTERNS: dict[str, tuple[str, ...]] = {
    "ONF": ("onf", "office national des forêts", "office national des forets"),
    "IGN": ("ign", "bd foret", "bd forêt", "lidar hd", "bd alti"),
    "BRGM": ("brgm", "bd geologique", "bd géologique", "geo050k"),
    "CNPF/CRPF": ("cnpf", "crpf", "centre regional de la propriete forestiere"),
    "INRAE": ("inrae", "irstea", "inra "),
    "Météo-France": ("meteo-france", "météo-france", "meteofrance"),
    "OFB/ONCFS": ("ofb", "oncfs"),
    "PNR Millevaches": ("pnr", "millevaches"),
    "OpenStreetMap": ("openstreetmap", " osm", "osm."),
    "BTS GF / production personnelle": (
        "camille perraudeau", "fiche terrain", "diagnostic stationnel camille",
        "projet fin d anner", "projet fin d'année",
    ),
}

SENSITIVE_PATTERNS = (
    "proprietaire", "propriétaire", "parcelles propriétaire", "parcelle privee",
    "parcelle privée", "coordonnees", "coordonnées", "telephone", "téléphone",
    "adresse", "contact", "cadastre", "psg", "plan simple de gestion",
)


def normalize(value: str) -> str:
    return re.sub(r"\s+", " ", value.casefold()).strip()


def safe_text(path: Path, max_chars: int = 20_000) -> tuple[str, str]:
    """Extrait un texte borné ; retourne aussi le statut d'extraction."""

    suffix = path.suffix.casefold()
    try:
        if suffix == ".pdf":
            reader = PdfReader(str(path), strict=False)
            pdf_chunks = [(page.extract_text() or "") for page in reader.pages[:8]]
            text = "\n".join(pdf_chunks)[:max_chars]
            status = "texte_extrait" if len(text.strip()) >= 120 else "ocr_probablement_requis"
            return text, f"{status};pages={len(reader.pages)}"
        if suffix == ".docx":
            document = Document(str(path))
            docx_chunks = [p.text for p in document.paragraphs]
            docx_chunks.extend(" | ".join(cell.text for cell in row.cells) for table in document.tables for row in table.rows)
            return "\n".join(docx_chunks)[:max_chars], "texte_extrait"
        if suffix == ".pptx":
            presentation = Presentation(str(path))
            pptx_chunks = [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")]
            return "\n".join(pptx_chunks)[:max_chars], f"texte_extrait;diapositives={len(presentation.slides)}"
        if suffix == ".odt":
            with zipfile.ZipFile(path) as archive:
                xml = archive.read("content.xml")
            root = ElementTree.fromstring(xml)
            return " ".join(root.itertext())[:max_chars], "texte_extrait"
        if suffix == ".xlsx":
            workbook = load_workbook(path, read_only=True, data_only=True)
            xlsx_chunks: list[str] = []
            for worksheet in workbook.worksheets[:8]:
                xlsx_chunks.append(worksheet.title)
                for row in worksheet.iter_rows(max_row=100, values_only=True):
                    xlsx_chunks.append(" | ".join(str(value) for value in row if value is not None))
                    if sum(map(len, xlsx_chunks)) >= max_chars:
                        break
            workbook.close()
            return "\n".join(xlsx_chunks)[:max_chars], "texte_extrait"
        if suffix in {".csv", ".tsv", ".txt", ".md", ".geojson", ".json"}:
            return path.read_text(encoding="utf-8", errors="replace")[:max_chars], "texte_extrait"
    except Exception as exc:  # les fichiers bureautiques historiques peuvent être corrompus
        return "", f"echec:{type(exc).__name__}"
    return "", "non_extrait"


def gpkg_metadata(path: Path) -> tuple[str, str]:
    try:
        with sqlite3.connect(f"file:{path.as_posix()}?mode=ro", uri=True) as connection:
            rows = connection.execute(
                "SELECT table_name, data_type, identifier, srs_id FROM gpkg_contents ORDER BY table_name"
            ).fetchall()
        text = "; ".join(f"{name}|{kind}|{identifier}|EPSG:{srs}" for name, kind, identifier, srs in rows)
        return text, f"geopackage_lisible;couches={len(rows)}"
    except Exception as exc:
        return "", f"geopackage_echec:{type(exc).__name__}"


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def top_root(path: Path, root: Path) -> str:
    relative = path.relative_to(root)
    return relative.parts[0] if relative.parts else "<racine>"


def classify(path: Path) -> str | None:
    suffix = path.suffix.casefold()
    if suffix in DOCUMENT_EXTENSIONS:
        return "document"
    if suffix in TABLE_EXTENSIONS:
        return "tableur_texte_structure"
    if suffix in RASTER_EXTENSIONS:
        return "raster_geospatial"
    if suffix == ".shp":
        return "vecteur_shapefile"
    if suffix in VECTOR_EXTENSIONS:
        return "vecteur_geospatial"
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in {".txt", ".md", ".json"}:
        return "texte"
    return None


def logical_files(files: list[Path]) -> tuple[list[dict[str, Any]], Counter[str]]:
    """Réunit les sidecars d'un shapefile en une ressource logique."""

    by_key = {str(path).casefold(): path for path in files}
    consumed: set[str] = set()
    resources: list[dict[str, Any]] = []
    sidecar_counts: Counter[str] = Counter()

    for path in sorted(files, key=lambda item: str(item).casefold()):
        key = str(path).casefold()
        if key in consumed:
            continue
        category = classify(path)
        if category is None:
            continue
        members = [path]
        if path.suffix.casefold() == ".shp":
            for suffix in SHAPEFILE_SIDECARS:
                candidate = path.with_suffix(suffix)
                candidate_key = str(candidate).casefold()
                if candidate_key in by_key:
                    members.append(by_key[candidate_key])
                    consumed.add(candidate_key)
                    sidecar_counts[suffix] += 1
        consumed.add(key)
        resources.append({"path": path, "category": category, "members": members})
    return resources, sidecar_counts


def assess(path: Path, category: str, extracted: str) -> dict[str, Any]:
    corpus = normalize(f"{path} {extracted}")
    normalized_path = normalize(str(path)).replace("/", "\\")
    excluded_reason = ""
    if any(token in normalized_path for token in EXCLUDED_PATH_TOKENS):
        excluded_reason = "artefact technique ou document personnel"
    domains = [domain for domain, patterns in DOMAIN_PATTERNS.items() if any(pattern in corpus for pattern in patterns)]
    score = min(100, len(domains) * 12)
    if "diagnostic_stationnel" in domains:
        score += 18
    if category in {"raster_geospatial", "vecteur_shapefile", "vecteur_geospatial"}:
        score += 10
    if any(token in corpus for token in ("anglais", "volley", "lettre de motivation", "conseil des delegues")):
        score -= 35
    if category == "image" and "geospatial_teledetection" not in domains and "incendie_dfci" not in domains:
        score = min(score, 12)
    if excluded_reason:
        score = 0
    score = max(0, min(100, score))

    if score >= 60:
        relevance = "élevée"
    elif score >= 30:
        relevance = "moyenne"
    elif score > 0:
        relevance = "faible"
    else:
        relevance = "hors_périmètre_probable"

    producers = [name for name, patterns in PRODUCER_PATTERNS.items() if any(pattern in corpus for pattern in patterns)]
    producer = "; ".join(producers) if producers else "à_identifier"

    sensitive = any(pattern in corpus for pattern in SENSITIVE_PATTERNS)
    sensitivity = "restreinte_à_confirmer" if sensitive else "non_détectée"

    if "ONF" in producers:
        rights = "citation_seule_probable;licence_et_reproduction_à_confirmer"
    elif any(item in producers for item in ("IGN", "BRGM", "PNR Millevaches", "CNPF/CRPF")):
        rights = "licence_source_et_millésime_à_confirmer"
    elif "OpenStreetMap" in producers:
        rights = "ODbL_probable;provenance_et_attribution_à_confirmer"
    elif "BTS GF / production personnelle" in producers:
        rights = "production_personnelle;coauteurs_et_données_tiers_à_confirmer"
    else:
        rights = "inconnus;ingestion_interdite"

    if excluded_reason:
        route = "EXCLURE_CONFIDENTIEL_OU_TECHNIQUE"
    elif "diagnostic_stationnel" in domains or "benchmark_formation" in domains:
        route = "GSIE_BENCH_et_FIELD_INTAKE"
    elif category in {"raster_geospatial", "vecteur_shapefile", "vecteur_geospatial"}:
        route = "DATA_REGISTRY_metadata_only"
    elif domains:
        route = "RESEARCH_puis_KNOWLEDGE_après_validation"
    else:
        route = "ARCHIVER_HORS_PÉRIMÈTRE"

    action = "QUALIFIER" if relevance in {"élevée", "moyenne"} else "TRIER_MANUELLEMENT"
    if relevance == "hors_périmètre_probable":
        action = "EXCLURE_PROVISOIREMENT"
    return {
        "score_pertinence": score,
        "pertinence": relevance,
        "domaines": domains,
        "producteur_apparent": producer,
        "droits": rights,
        "sensibilite": sensitivity,
        "routage": route,
        "action": action,
        "motif_exclusion": excluded_reason,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", type=Path, default=Path(r"E:\Documents"))
    parser.add_argument("--output", type=Path, default=Path(__file__).parent / "inventory_edocuments")
    parser.add_argument("--hash-max-mib", type=int, default=256)
    args = parser.parse_args()

    source = args.source.resolve()
    output = args.output.resolve()
    output.mkdir(parents=True, exist_ok=True)

    all_files = [path for path in source.rglob("*") if path.is_file()]
    excluded_roots: Counter[str] = Counter()
    excluded_paths: Counter[str] = Counter()
    eligible: list[Path] = []
    for candidate_path in all_files:
        root_name = top_root(candidate_path, source)
        if root_name.casefold() in EXCLUDED_ROOTS:
            excluded_roots[root_name] += 1
            continue
        normalized_path = normalize(str(candidate_path)).replace("/", "\\")
        matching_token = next((token for token in EXCLUDED_PATH_TOKENS if token in normalized_path), None)
        if matching_token:
            excluded_paths[matching_token] += 1
            continue
        eligible.append(candidate_path)

    resources, sidecar_counts = logical_files(eligible)
    rows: list[dict[str, Any]] = []
    for index, resource in enumerate(resources, start=1):
        path: Path = resource["path"]
        category: str = resource["category"]
        members: list[Path] = resource["members"]
        size = sum(member.stat().st_size for member in members)
        extracted = ""
        extraction_status = "non_extrait"
        if category in {"document", "tableur_texte_structure", "texte"} and size <= 100 * 1024 * 1024:
            extracted, extraction_status = safe_text(path)
        elif path.suffix.casefold() == ".gpkg":
            extracted, extraction_status = gpkg_metadata(path)
        elif path.suffix.casefold() in {".geojson", ".json"} and size <= 20 * 1024 * 1024:
            extracted, extraction_status = safe_text(path)

        assessment = assess(path, category, extracted)
        checksum = ""
        checksum_status = "non_calculé"
        if size <= args.hash_max_mib * 1024 * 1024 and assessment["pertinence"] != "hors_périmètre_probable":
            checksum = sha256(path)
            checksum_status = "sha256_fichier_principal"
        elif size > args.hash_max_mib * 1024 * 1024:
            checksum_status = "différé_taille_supérieure_à_la_borne"

        rows.append(
            {
                "resource_id": f"EDOC-{index:05d}",
                "chemin": str(path),
                "racine": top_root(path, source),
                "categorie": category,
                "format": path.suffix.casefold().lstrip("."),
                "taille_octets_logique": size,
                "membres_logiques": len(members),
                "fichiers_associes": " | ".join(str(member) for member in members[1:]),
                "modifie_le": datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec="seconds"),
                "sha256": checksum,
                "statut_checksum": checksum_status,
                "statut_extraction": extraction_status,
                **assessment,
            }
        )

    duplicate_map: dict[str, list[str]] = defaultdict(list)
    for row in rows:
        if row["sha256"]:
            duplicate_map[row["sha256"]].append(row["resource_id"])
    duplicate_groups = {digest: ids for digest, ids in duplicate_map.items() if len(ids) > 1}
    group_by_id: dict[str, str] = {}
    for number, ids in enumerate(sorted(duplicate_groups.values(), key=lambda values: values[0]), start=1):
        for resource_id in ids:
            group_by_id[resource_id] = f"DUP-{number:04d}"
    for row in rows:
        row["groupe_doublon"] = group_by_id.get(row["resource_id"], "")
        row["domaines"] = ";".join(row["domaines"])

    rows.sort(key=lambda row: (-int(row["score_pertinence"]), row["chemin"].casefold()))
    fields = list(rows[0]) if rows else []
    with (output / "manifest.csv").open("w", encoding="utf-8-sig", newline="") as stream:
        writer = csv.DictWriter(stream, fieldnames=fields)
        writer.writeheader()
        writer.writerows(rows)
    (output / "manifest.json").write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")

    summary = {
        "generated_at": datetime.now().astimezone().isoformat(timespec="seconds"),
        "source": str(source),
        "files_seen": len(all_files),
        "eligible_files_after_root_exclusions": len(eligible),
        "logical_resources": len(rows),
        "categories": Counter(row["categorie"] for row in rows),
        "relevance": Counter(row["pertinence"] for row in rows),
        "routes": Counter(row["routage"] for row in rows),
        "rights": Counter(row["droits"] for row in rows),
        "sensitivity": Counter(row["sensibilite"] for row in rows),
        "extraction_status": Counter(row["statut_extraction"].split(";")[0] for row in rows),
        "excluded_roots": excluded_roots,
        "excluded_path_patterns": excluded_paths,
        "shapefile_sidecars_folded": sidecar_counts,
        "duplicate_groups": len(duplicate_groups),
        "duplicate_resources": sum(len(ids) for ids in duplicate_groups.values()),
        "checksum_scope": f"fichier principal <= {args.hash_max_mib} MiB et pertinence non nulle",
    }
    (output / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
